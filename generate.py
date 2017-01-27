# -*- coding: utf-8 -*-

import csv
import sys
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
import copy
import six


def _get_blank_slide_layout(pres):
	layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
	min_items = min(layout_items_count)
	blank_layout_id = layout_items_count.index(min_items)
	return pres.slide_layouts[blank_layout_id]

def copy_slide(pres,pres1,index):
	source = pres.slides[index]
	blank_slide_layout = _get_blank_slide_layout(pres)
	dest = pres1.slides.add_slide(blank_slide_layout)
	
    
	for shp in source.shapes:
		el = shp.element
		newel = copy.deepcopy(el)
		dest.shapes._spTree.insert_element_before(newel, 'p:extLst')
		
		for key, value in six.iteritems(source.part.rels):
			if not "notesSlide" in value.reltype:
				dest.rels.add_relationship(value.reltype, value._target, value.rId)

		return dest

		
def duplicate_slide(pres, index):
	"""Duplicate the slide with the given index in pres.

    Adds slide to the end of the presentation"""
	source = pres.slides[index]
	blank_slide_layout = _get_blank_slide_layout(pres)
	dest = pres.slides.add_slide(blank_slide_layout)
	for shp in source.shapes:
		el = shp.element
		newel = copy.deepcopy(el)
		dest.shapes._spTree.insert_element_before(newel, 'p:extLst')
		
	for key, value in six.iteritems(source.part.rels):
	# Make sure we don't copy a notesSlide relation as that won't exist
		if not "notesSlide" in value.reltype:
			dest.rels.add_relationship(value.reltype, value._target, value.rId)
	return dest
	
	
prs=Presentation("/root/vuln.pptx")
f = open(sys.argv[1], 'rb')
reader = csv.DictReader(f,delimiter=';')
i=0
for row in reader:
 ref= row['EY reference'].decode('iso-8859-1').encode('utf8')
 title= row['Titre'].decode('iso-8859-1').encode('utf8')
 description= row['Description'].decode('iso-8859-1').encode('utf8')
 impact= row['Impact'].decode('iso-8859-1').encode('utf8')
 recommendation= row['Recommandation'].decode('iso-8859-1').encode('utf8')
 severity = row['Severite'].decode('iso-8859-1').encode('utf8')
 host = row['host'].decode('iso-8859-1').encode('utf8')
 
 
 slide=prs.slides[i]
 table = slide.shapes[5].table
 
 
 #adding reference
 table.rows[0].cells[0].text_frame.text = ref
 table.rows[0].cells[0].text_frame.paragraphs[0].font.name="EYInterstate Light"
 table.rows[0].cells[0].text_frame.paragraphs[0].font.size=Pt(10)
 table.rows[0].cells[0].text_frame.paragraphs[0].font.bold = False
 table.rows[0].cells[0].text_frame.paragraphs[0].font.italic = False
 
 #adding title
 table.rows[0].cells[1].text_frame.text = title
 table.rows[0].cells[1].text_frame.paragraphs[0].font.name="EYInterstate Light"
 table.rows[0].cells[1].text_frame.paragraphs[0].font.size=Pt(10)
 table.rows[0].cells[1].text_frame.paragraphs[0].font.bold = True
 table.rows[0].cells[1].text_frame.paragraphs[0].font.italic = False
 
 #adding host
 table.rows[1].cells[1].text_frame.text = host
 table.rows[1].cells[1].text_frame.paragraphs[0].font.name="EYInterstate Light"
 table.rows[1].cells[1].text_frame.paragraphs[0].font.size=Pt(10)
 table.rows[1].cells[1].text_frame.paragraphs[0].font.bold = False
 table.rows[1].cells[1].text_frame.paragraphs[0].font.italic = False
 
 #adding description
 table.rows[3].cells[0].text_frame.text = description
 table.rows[3].cells[0].text_frame.paragraphs[0].font.name="EYInterstate Light"
 table.rows[3].cells[0].text_frame.paragraphs[0].font.size=Pt(10)
 table.rows[3].cells[0].text_frame.paragraphs[0].font.bold = False
 table.rows[3].cells[0].text_frame.paragraphs[0].font.italic = True
 table.rows[3].cells[0].text_frame.paragraphs[0].font.color.rgb = RGBColor(60,60,60)
 
 #adding impact
 table.rows[5].cells[0].text_frame.text = impact
 table.rows[5].cells[0].text_frame.paragraphs[0].font.name="EYInterstate Light"
 table.rows[5].cells[0].text_frame.paragraphs[0].font.size=Pt(10)
 table.rows[5].cells[0].text_frame.paragraphs[0].font.bold = False
 table.rows[5].cells[0].text_frame.paragraphs[0].font.italic = True
 table.rows[5].cells[0].text_frame.paragraphs[0].font.color.rgb = RGBColor(60,60,60)
 
 #adding recommendation
 table.rows[7].cells[0].text_frame.text = recommendation
 table.rows[7].cells[0].text_frame.paragraphs[0].font.name="EYInterstate Light"
 table.rows[7].cells[0].text_frame.paragraphs[0].font.size=Pt(10)
 table.rows[7].cells[0].text_frame.paragraphs[0].font.bold = False
 table.rows[7].cells[0].text_frame.paragraphs[0].font.italic = True
 table.rows[7].cells[0].text_frame.paragraphs[0].font.color.rgb = RGBColor(60,60,60)
 
 #adding severity
 table.rows[0].cells[3].text_frame.text = severity
 table.rows[0].cells[3].text_frame.paragraphs[0].font.name="EYInterstate Light"
 table.rows[0].cells[3].text_frame.paragraphs[0].font.size=Pt(10)
 table.rows[0].cells[3].text_frame.paragraphs[0].font.bold = False
 table.rows[0].cells[3].text_frame.paragraphs[0].font.italic = True
 table.rows[0].cells[3].text_frame.paragraphs[0].font.color.rgb = RGBColor(60,60,60)
 table.rows[0].cells[3].fill.solid()
 
 if severity=="D-Critical":
	table.rows[0].cells[3].fill.fore_color.rgb = RGBColor(192, 0, 0)
 if severity=="C-High":
	table.rows[0].cells[3].fill.fore_color.rgb = RGBColor(255, 158, 0)
 if severity=="B-Medium":
	table.rows[0].cells[3].fill.fore_color.rgb = RGBColor(255, 204, 0)
 if severity=="A-Low":
	table.rows[0].cells[3].fill.fore_color.rgb = RGBColor(255, 250, 204)
 
 i=i+1
 
 
prs.save("/var/www/html/findings.pptx")
